"""
Export Union Eyes buyer deck (UNION_EYES_BUYER_DECK.md) to PPTX.

Usage:
    python scripts/export_buyer_deck.py
    python scripts/export_buyer_deck.py --out demo-output/buyer-deck.pptx

Requires: python-pptx (pip install python-pptx)
"""

from __future__ import annotations

import argparse
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path

# ── Colours (Union Eyes brand: dark navy + white + accent teal) ───────────────
NAVY = "1B2A4A"  # primary background / title bg
WHITE = "FFFFFF"
TEAL = "00A896"  # accent bullets / line
LIGHT = "F4F6F9"  # body background on content slides
GREY = "6B7A99"  # speaker-notes meta text (unused in slide itself)

DECK_MD = (
    Path(__file__).parent.parent
    / "docs/commercial/close-package/UNION_EYES_BUYER_DECK.md"
)


# ── Slide data structure ──────────────────────────────────────────────────────


@dataclass
class Slide:
    number: int  # 0 = cover/preamble, 1–13 = main, 100+ = appendix
    label: str  # e.g. "Slide 1 — Title / Why Now"
    title: str  # **Title:** value
    bullets: list[str] = field(default_factory=list)
    notes: str = ""
    is_appendix: bool = False
    is_section_break: bool = False  # used for "Appendix" divider


# ── Parser ────────────────────────────────────────────────────────────────────

_SLIDE_RE = re.compile(r"^## (Slide (\d+) — .+)$")
_APPENDIX_RE = re.compile(r"^## (Appendix slides)")
_TITLE_RE = re.compile(r"^\*\*Title:\*\*\s+(.+)$")
_BULLET_RE = re.compile(r"^[-*]\s+(.+)$")
_NOTES_RE = re.compile(r"^\*\*Speaker notes?:\*\*\s*(.*)$")


def _strip_md(text: str) -> str:
    """Remove bold/italic markers for plain text in PPTX."""
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    # Remove inline links [text](url) → text
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    return text.strip()


def parse_deck(md_path: Path) -> list[Slide]:
    slides: list[Slide] = []
    current: Slide | None = None
    in_bullets = False
    in_notes = False
    notes_lines: list[str] = []

    def _flush(s: Slide | None, nl: list[str]) -> None:
        if s:
            s.notes = _strip_md(" ".join(nl).strip())

    for raw in md_path.read_text(encoding="utf-8").splitlines():
        line = raw.strip()

        m = _SLIDE_RE.match(line)
        if m:
            _flush(current, notes_lines)
            notes_lines = []
            in_bullets = in_notes = False
            num = int(m.group(2))
            current = Slide(number=num, label=_strip_md(m.group(1)), title="")
            slides.append(current)
            continue

        if _APPENDIX_RE.match(line):
            _flush(current, notes_lines)
            notes_lines = []
            in_bullets = in_notes = False
            # Insert a section-break pseudo slide
            current = Slide(
                number=100,
                label="Appendix",
                title="Appendix Slides",
                is_section_break=True,
                is_appendix=True,
            )
            slides.append(current)
            continue

        if current is None:
            continue

        m = _TITLE_RE.match(line)
        if m:
            current.title = _strip_md(m.group(1))
            in_bullets = in_notes = False
            continue

        if line == "**Bullets:**":
            in_bullets = True
            in_notes = False
            continue

        m = _NOTES_RE.match(line)
        if m:
            in_bullets = False
            in_notes = True
            first = _strip_md(m.group(1))
            if first:
                notes_lines.append(first)
            continue

        if in_bullets:
            m = _BULLET_RE.match(line)
            if m:
                current.bullets.append(_strip_md(m.group(1)))
            elif line == "" or line.startswith("---"):
                in_bullets = False
            continue

        if in_notes:
            if line == "" or line.startswith("---") or line.startswith("##"):
                # end of notes block — don't stop on blank between sentences
                if line.startswith("---") or line.startswith("##"):
                    in_notes = False
                continue
            notes_lines.append(_strip_md(line))

    _flush(current, notes_lines)
    return slides


# ── PPTX builder ─────────────────────────────────────────────────────────────

try:
    import lxml.etree as etree
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Emu, Inches, Pt
except ImportError:
    sys.exit("python-pptx not found. Run: pip install python-pptx")


def _rgb(hex6: str) -> RGBColor:
    r, g, b = int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16)
    return RGBColor(r, g, b)


def _fill_solid(shape, hex6: str) -> None:
    from pptx.dml.color import RGBColor

    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex6)


def _add_textbox(
    slide,
    left_in,
    top_in,
    w_in,
    h_in,
    text: str,
    font_size: int,
    bold=False,
    color=WHITE,
    align=PP_ALIGN.LEFT,
    wrap=True,
) -> None:
    from pptx.util import Inches, Pt

    txb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(w_in), Inches(h_in)
    )
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


def _set_slide_background(slide, hex6: str) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex6)


def _add_notes(slide, text: str) -> None:
    if not text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def _build_cover_slide(prs: Presentation, slide_data: Slide) -> None:
    """Slide 1: full navy background, large title, subtitle line."""
    blank_layout = prs.slide_layouts[6]  # blank
    sl = prs.slides.add_slide(blank_layout)
    _set_slide_background(sl, NAVY)

    # Teal accent bar (left edge)
    bar = sl.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(0.12), Inches(7.5)  # MSO_SHAPE_TYPE.RECTANGLE
    )
    _fill_solid(bar, TEAL)
    bar.line.fill.background()

    # Slide number badge
    _add_textbox(sl, 0.25, 0.25, 0.8, 0.35, "01", 11, bold=True, color=TEAL)

    # Main title
    txb = sl.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(8.8), Inches(2.0))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = slide_data.title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = _rgb(WHITE)

    # Bullets as subtitle taglines
    if slide_data.bullets:
        y = 4.0
        for b in slide_data.bullets:
            txb2 = sl.shapes.add_textbox(
                Inches(0.6), Inches(y), Inches(8.8), Inches(0.45)
            )
            txb2.word_wrap = True
            tf2 = txb2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            run2 = p2.add_run()
            run2.text = "  " + b
            run2.font.size = Pt(16)
            run2.font.color.rgb = _rgb(LIGHT)
            y += 0.46

    # Footer
    _add_textbox(
        sl,
        0.6,
        6.9,
        8.8,
        0.4,
        "Union Eyes  ·  Confidential  ·  April 2026",
        9,
        color=GREY,
    )

    _add_notes(sl, slide_data.notes)


def _build_section_break(prs: Presentation, title: str = "Appendix") -> None:
    blank_layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(blank_layout)
    _set_slide_background(sl, NAVY)
    bar = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    _fill_solid(bar, TEAL)
    bar.line.fill.background()
    _add_textbox(
        sl, 0.6, 3.0, 8.8, 1.2, title, 40, bold=True, color=WHITE, align=PP_ALIGN.CENTER
    )


def _build_content_slide(prs: Presentation, slide_data: Slide, slide_num: int) -> None:
    """Standard content slide: light background, navy header bar, teal bullet dots."""
    blank_layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(blank_layout)
    _set_slide_background(sl, LIGHT)

    # Navy header bar
    header = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.35))
    _fill_solid(header, NAVY)
    header.line.fill.background()

    # Slide number in header
    num_str = f"{slide_num:02d}"
    _add_textbox(sl, 0.2, 0.05, 0.55, 0.55, num_str, 11, bold=True, color=TEAL)

    # Title in header
    txb = sl.shapes.add_textbox(Inches(0.75), Inches(0.08), Inches(8.8), Inches(1.0))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = slide_data.title
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = _rgb(WHITE)

    # Bullets
    if slide_data.bullets:
        y = 1.55
        row_h = min(0.72, (7.5 - 1.8) / max(len(slide_data.bullets), 1))
        font_size = 15 if len(slide_data.bullets) <= 6 else 13

        for b in slide_data.bullets:
            # Teal dot
            dot = sl.shapes.add_shape(
                1, Inches(0.35), Inches(y + 0.12), Inches(0.12), Inches(0.12)
            )
            _fill_solid(dot, TEAL)
            dot.line.fill.background()

            # Bullet text
            txb2 = sl.shapes.add_textbox(
                Inches(0.6), Inches(y), Inches(9.1), Inches(row_h)
            )
            txb2.word_wrap = True
            tf2 = txb2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            run2 = p2.add_run()
            run2.text = b
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = _rgb(NAVY)
            y += row_h

    # Footer bar
    footer = sl.shapes.add_shape(1, Inches(0), Inches(7.1), Inches(10), Inches(0.4))
    _fill_solid(footer, NAVY)
    footer.line.fill.background()
    _add_textbox(
        sl,
        0.25,
        7.12,
        9.5,
        0.3,
        "Union Eyes  ·  Confidential  ·  unioneyes.ca",
        8,
        color=GREY,
    )

    _add_notes(sl, slide_data.notes)


def build_pptx(slides: list[Slide], out_path: Path) -> None:
    # Widescreen 16:9
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for s in slides:
        if s.is_section_break:
            _build_section_break(prs, s.title)
        elif s.number == 1:
            _build_cover_slide(prs, s)
        else:
            _build_content_slide(prs, s, s.number)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"✓  Saved PPTX → {out_path}")


# ── Main ──────────────────────────────────────────────────────────────────────


def main() -> None:
    parser = argparse.ArgumentParser(description="Export Union Eyes buyer deck to PPTX")
    parser.add_argument(
        "--deck", default=str(DECK_MD), help="Path to UNION_EYES_BUYER_DECK.md"
    )
    parser.add_argument(
        "--out",
        default="demo-output/union-eyes-buyer-deck.pptx",
        help="Output PPTX path",
    )
    args = parser.parse_args()

    md_path = Path(args.deck)
    out_path = Path(args.out)

    if not md_path.exists():
        sys.exit(f"Deck not found: {md_path}")

    print(f"Parsing  {md_path}")
    slides = parse_deck(md_path)
    print(f"Found {len(slides)} slides/sections")

    build_pptx(slides, out_path)

    size_kb = out_path.stat().st_size // 1024
    print(f"   Size: {size_kb} KB")
    print()
    print("To open:")
    print(f"  start {out_path}")
    print()
    print(
        "Note: PDF export requires LibreOffice or MS Office. See scripts/Export-BuyerDeck.ps1 for --pdf flag."
    )


if __name__ == "__main__":
    main()
